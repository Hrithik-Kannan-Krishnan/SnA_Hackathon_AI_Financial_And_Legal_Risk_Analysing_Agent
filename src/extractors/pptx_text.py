from __future__ import annotations
from pptx import Presentation

def extract_pptx_text(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    slide_text.append(t)
        if slide_text:
            parts.append(f"Slide {i}:\n" + "\n".join(slide_text))
    return "\n\n".join(parts)
